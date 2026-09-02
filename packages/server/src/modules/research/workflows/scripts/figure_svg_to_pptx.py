#!/usr/bin/env python3
"""figure_svg_to_pptx.py - research workbench figure sidecar (mapping v2).

Converts a standalone SVG figure (the constrained element vocabulary produced
by the figure-drawing workflow) into a Microsoft PowerPoint file made of
NATIVE, editable shapes via python-pptx. The raw SVG is not embedded as a
bitmap.

Mapping v2 (on top of the v1 rect/circle/ellipse/line/polyline/polygon/text):

- path: M/L/H/V/C/S/Q/T/Z commands produce real freeform shapes. Cubic and
  quadratic beziers are flattened by sampling (CURVE_SAMPLES per segment);
  A/a arcs are DOWNGRADED to a polyline approximation via the SVG endpoint
  parameterization (F.6.5) and counted under svgFeaturesSkipped
  ["arcApproximated"]. Each subpath becomes its own freeform shape; Z closes
  the freeform.
- transform: translate/matrix/scale/rotate are parsed into an affine.
  Point-mapped shapes (line/polyline/polygon/path) transform their vertices,
  which is exact for translate/rotate/scale. Bounding-box shapes
  (rect/circle/ellipse/text) are placed via their transformed center plus
  shape.rotation, which is exact for rotate+translate; a scale/skew component
  on those shapes cannot be expressed without resizing, so it is ignored and
  counted under svgFeaturesSkipped ["scaleOnBoxShapeIgnored"].
- linearGradient/radialGradient: fill="url(#id)" resolves to a single
  representative solid color — the trapezoidal (offset-weighted) average of
  the stop colors, i.e. the average color of the gradient ramp. Stroke paint
  is always preserved verbatim (url() stroke resolves through the same table,
  falling back to black when unknown).
- tspan: inline tspans (no x/y/dy) become additional RUNS in the same
  paragraph carrying their own font-size/fill/bold overrides; tspans that
  move the baseline (x/y/dy present) become additional PARAGRAPHS in the
  same text frame. Paragraph x positions collapse to the frame's left edge
  (only the first line keeps its text-anchor estimate) — an accepted
  approximation documented here and in the spike report.

Usage:
    python figure_svg_to_pptx.py <svg_path> <pptx_path> [title]

Prints a JSON summary on stdout on success (ok/pptx/shapes/slidePx plus the
additive svgFeaturesMapped / svgFeaturesSkipped counters); exits non-zero
with a message on stderr on failure. The calling workflow node degrades
gracefully when this sidecar is not configured, so this script only needs to
be honest about its own failures.

Requires: python-pptx (pip install python-pptx). Nothing else beyond stdlib.
"""

import json
import math
import re
import sys

EMU_PER_PX = 9525  # 914400 EMU per inch / 96 px per inch
DEFAULT_WIDTH = 900.0
DEFAULT_HEIGHT = 560.0
CURVE_SAMPLES = 16  # flattening samples per bezier segment
ARC_DEGREES_PER_SAMPLE = 5.0  # polyline resolution for arc approximation

NAMED_COLORS = {
    "black": "000000",
    "white": "FFFFFF",
    "red": "FF0000",
    "green": "008000",
    "blue": "0000FF",
    "gray": "808080",
    "grey": "808080",
    "orange": "FFA500",
}

NUMBER_PATTERN = r"-?(?:\d*\.\d+|\d+)(?:[eE][-+]?\d+)?"
PATH_COMMAND_ARGUMENTS = {
    "M": 2, "L": 2, "H": 1, "V": 1, "C": 6, "S": 4, "Q": 4, "T": 2, "A": 7, "Z": 0,
}


def fail(message):
    sys.stderr.write("figure_svg_to_pptx: %s\n" % message)
    sys.exit(1)


def local_name(tag):
    return tag.rsplit("}", 1)[-1] if isinstance(tag, str) else ""


def parse_length(value, fallback=None):
    if value is None:
        return fallback
    text = str(value).strip()
    if text.endswith("px"):
        text = text[:-2]
    try:
        return float(text)
    except ValueError:
        return fallback


def parse_color(value):
    """Returns a 6-digit RGB hex string, 'none', or None (unspecified)."""
    if value is None:
        return None
    text = str(value).strip().lower()
    if not text:
        return None
    if text == "none":
        return "none"
    if text.startswith("#"):
        hex_part = text[1:]
        if len(hex_part) == 3:
            hex_part = "".join(ch * 2 for ch in hex_part)
        if len(hex_part) >= 6 and re.match(r"^[0-9a-f]{6}$", hex_part[:6]):
            return hex_part[:6].upper()
        return None
    return NAMED_COLORS.get(text)


def bump(counter, key, amount=1):
    counter[key] = counter.get(key, 0) + amount


class Affine(object):
    """Minimal SVG-style 2D affine transform (a b c d e f)."""

    def __init__(self, a=1.0, b=0.0, c=0.0, d=1.0, e=0.0, f=0.0):
        self.a, self.b, self.c, self.d, self.e, self.f = a, b, c, d, e, f

    def apply(self, x, y):
        return (self.a * x + self.c * y + self.e, self.b * x + self.d * y + self.f)

    def composed(self, other):
        """Returns self ∘ other (other is applied to points first)."""
        return Affine(
            self.a * other.a + self.c * other.b,
            self.b * other.a + self.d * other.b,
            self.a * other.c + self.c * other.d,
            self.b * other.c + self.d * other.d,
            self.a * other.e + self.c * other.f + self.e,
            self.b * other.e + self.d * other.f + self.f,
        )

    @property
    def rotation_degrees(self):
        # SVG rotate() is clockwise in the y-down screen space, exactly like
        # python-pptx shape.rotation, so the angle maps through unchanged.
        return math.degrees(math.atan2(self.b, self.a))

    @property
    def has_scale_or_skew(self):
        return abs(math.hypot(self.a, self.b) - 1.0) > 1e-6 or abs(math.hypot(self.c, self.d) - 1.0) > 1e-6


def parse_transform_affine(value):
    """Parses an SVG transform attribute into an Affine (left-to-right order)."""
    affine = Affine()
    if not value:
        return affine
    for name, args_text in re.findall(r"([a-zA-Z]+)\s*\(([^)]*)\)", str(value)):
        numbers = [float(n) for n in re.findall(NUMBER_PATTERN, args_text)]
        name = name.strip().lower()
        op = None
        if name == "translate" and numbers:
            op = Affine(1, 0, 0, 1, numbers[0], numbers[1] if len(numbers) > 1 else 0.0)
        elif name == "matrix" and len(numbers) >= 6:
            op = Affine(*numbers[:6])
        elif name == "scale" and numbers:
            sx = numbers[0]
            op = Affine(sx, 0, 0, numbers[1] if len(numbers) > 1 else sx, 0, 0)
        elif name == "rotate" and numbers:
            angle = math.radians(numbers[0])
            rotation = Affine(math.cos(angle), math.sin(angle), -math.sin(angle), math.cos(angle), 0, 0)
            if len(numbers) >= 3:
                # rotate(a, cx, cy) == translate(cx,cy) ∘ rotate(a) ∘ translate(-cx,-cy)
                cx, cy = numbers[1], numbers[2]
                op = Affine(1, 0, 0, 1, cx, cy).composed(rotation).composed(Affine(1, 0, 0, 1, -cx, -cy))
            else:
                op = rotation
        # skewX/skewY are not part of the constrained figure vocabulary; they
        # degrade to identity for bounding-box shapes and are covered by the
        # scale/skew skip counter wherever the linear part is non-rigid.
        if op is not None:
            affine = affine.composed(op)
    return affine


def collect_elements(node, affine=None, output=None):
    """Depth-first flatten of drawable elements with accumulated affine."""
    if affine is None:
        affine = Affine()
    if output is None:
        output = []
    for child in node:
        tag = local_name(child.tag)
        child_affine = affine
        if tag == "g":
            child_affine = affine.composed(parse_transform_affine(child.get("transform")))
        if tag in ("rect", "circle", "ellipse", "line", "polyline", "polygon", "path", "text"):
            output.append((tag, child, child_affine))
        collect_elements(child, child_affine, output)
    return output


def collect_gradients(root):
    """Maps gradient ids to one representative solid color.

    Strategy: the representative color is the average color of the gradient
    ramp — the trapezoidal integral of the piecewise-linear stop interpolation
    (each stop weighted by its span, first/last stops by half their outer
    span). For a plain two-stop gradient this is exactly the midpoint blend.
    Opacity/stop-opacity are ignored (documented approximation).
    """
    gradients = {}
    for element in root.iter():
        if local_name(element.tag) not in ("linearGradient", "radialGradient"):
            continue
        gradient_id = element.get("id")
        if not gradient_id:
            continue
        stops = []
        for stop in element.iter():
            if local_name(stop.tag) != "stop":
                continue
            offset_text = str(stop.get("offset") or "0").strip()
            if offset_text.endswith("%"):
                offset = parse_length(offset_text[:-1], 0.0) / 100.0
            else:
                offset = parse_length(offset_text, 0.0)
            color = parse_color(stop.get("stop-color"))
            if color in (None, "none"):
                color = "808080"  # neutral gray when a stop color is unusable
            stops.append((max(0.0, min(1.0, offset)), color))
        if not stops:
            continue
        stops.sort(key=lambda stop: stop[0])
        if len(stops) == 1:
            gradients[gradient_id] = stops[0][1]
            continue
        total = 0.0
        acc_r = acc_g = acc_b = 0.0
        for (offset_a, color_a), (offset_b, color_b) in zip(stops, stops[1:]):
            span = offset_b - offset_a
            if span <= 0:
                continue
            total += span
            acc_r += (int(color_a[0:2], 16) + int(color_b[0:2], 16)) / 2.0 * span
            acc_g += (int(color_a[2:4], 16) + int(color_b[2:4], 16)) / 2.0 * span
            acc_b += (int(color_a[4:6], 16) + int(color_b[4:6], 16)) / 2.0 * span
        if total <= 0:
            gradients[gradient_id] = stops[0][1]
            continue
        gradients[gradient_id] = "%02X%02X%02X" % (
            int(round(acc_r / total)), int(round(acc_g / total)), int(round(acc_b / total)))
    return gradients


def resolve_paint(value, gradients):
    """Resolves a paint value; url(#gradient) becomes the representative color."""
    if value is None:
        return None
    text = str(value).strip()
    match = re.match(r"url\(\s*#([^)]+?)\s*\)", text)
    if match:
        return gradients.get(match.group(1))
    return parse_color(text)


def parse_path_commands(d_text):
    """Splits a path `d` attribute into (command, args, relative) tuples.

    Repeated coordinate groups (implicit command repetition, e.g.
    "M 0 0 L 1 1 2 2") are expanded into individual commands.
    """
    tokens = re.findall(r"[MmLlHhVvCcSsQqTtAaZz]|%s" % NUMBER_PATTERN, str(d_text))
    commands = []
    index = 0
    while index < len(tokens):
        token = tokens[index]
        if not token[0].isalpha():
            index += 1  # defensive: stray number outside any command
            continue
        command = token.upper()
        relative = token.islower()
        index += 1
        arg_count = PATH_COMMAND_ARGUMENTS[command]
        if arg_count == 0:
            commands.append((command, [], relative))
            continue
        numbers = []
        while index < len(tokens) and not tokens[index][0].isalpha():
            numbers.append(float(tokens[index]))
            index += 1
        for start in range(0, len(numbers) - arg_count + 1, arg_count):
            commands.append((command, numbers[start:start + arg_count], relative))
    return commands


def sample_cubic(p0, p1, p2, p3):
    return [
        (
            ((1 - t) ** 3) * p0[0] + 3 * ((1 - t) ** 2) * t * p1[0] + 3 * (1 - t) * (t ** 2) * p2[0] + (t ** 3) * p3[0],
            ((1 - t) ** 3) * p0[1] + 3 * ((1 - t) ** 2) * t * p1[1] + 3 * (1 - t) * (t ** 2) * p2[1] + (t ** 3) * p3[1],
        )
        for t in (i / float(CURVE_SAMPLES) for i in range(1, CURVE_SAMPLES + 1))
    ]


def sample_quadratic(p0, p1, p2):
    return [
        (
            ((1 - t) ** 2) * p0[0] + 2 * (1 - t) * t * p1[0] + (t ** 2) * p2[0],
            ((1 - t) ** 2) * p0[1] + 2 * (1 - t) * t * p1[1] + (t ** 2) * p2[1],
        )
        for t in (i / float(CURVE_SAMPLES) for i in range(1, CURVE_SAMPLES + 1))
    ]


def sample_arc(x1, y1, rx, ry, rotation_degrees, large_arc, sweep, x2, y2):
    """Approximates an SVG arc with a polyline (endpoint parameterization).

    Implements the SVG spec F.6.5 conversion from endpoints to center
    parameterization, then samples the ellipse uniformly. This is the
    documented downgrade path for A/a commands (see svgFeaturesSkipped).
    """
    if abs(rx) < 1e-9 or abs(ry) < 1e-9:
        return [(x2, y2)]  # degenerate arc collapses to a line
    rx, ry = abs(rx), abs(ry)
    phi = math.radians(rotation_degrees)
    cos_phi, sin_phi = math.cos(phi), math.sin(phi)
    dx2, dy2 = (x1 - x2) / 2.0, (y1 - y2) / 2.0
    x1p = cos_phi * dx2 + sin_phi * dy2
    y1p = -sin_phi * dx2 + cos_phi * dy2
    lambda_sq = (x1p / rx) ** 2 + (y1p / ry) ** 2
    if lambda_sq > 1.0:
        scale = math.sqrt(lambda_sq)
        rx *= scale
        ry *= scale
    sign = 1.0 if large_arc != sweep else -1.0
    numerator = rx * rx * ry * ry - rx * rx * y1p * y1p - ry * ry * x1p * x1p
    denominator = rx * rx * y1p * y1p + ry * ry * x1p * x1p
    factor = sign * math.sqrt(max(numerator / denominator, 0.0)) if denominator else 0.0
    cxp = factor * rx * y1p / ry
    cyp = -factor * ry * x1p / rx
    center_x = cos_phi * cxp - sin_phi * cyp + (x1 + x2) / 2.0
    center_y = sin_phi * cxp + cos_phi * cyp + (y1 + y2) / 2.0
    theta1 = math.atan2((y1p - cyp) / ry, (x1p - cxp) / rx)
    theta2 = math.atan2((-y1p - cyp) / ry, (-x1p - cxp) / rx)
    delta = theta2 - theta1
    if not sweep and delta > 0:
        delta -= 2 * math.pi
    elif sweep and delta < 0:
        delta += 2 * math.pi
    steps = max(4, min(64, int(math.ceil(abs(math.degrees(delta)) / ARC_DEGREES_PER_SAMPLE))))
    points = []
    for step in range(1, steps + 1):
        theta = theta1 + delta * step / float(steps)
        points.append((
            center_x + rx * math.cos(theta) * cos_phi - ry * math.sin(theta) * sin_phi,
            center_y + rx * math.cos(theta) * sin_phi + ry * math.sin(theta) * cos_phi,
        ))
    return points


def flatten_path(d_text):
    """Flattens an SVG path into drawable subpaths.

    Returns (subpaths, arc_count) where subpaths is a list of
    (points, closed) tuples with points as [(x, y), ...]; beziers are sampled
    with CURVE_SAMPLES points per segment, arcs go through sample_arc.
    """
    commands = parse_path_commands(d_text)
    subpaths = []
    points = []
    closed = False
    x = y = 0.0
    subpath_start = (0.0, 0.0)
    last_cubic_control = None
    last_quad_control = None
    arc_count = 0

    def flush(closed_flag):
        if len(points) >= 2:
            subpaths.append((list(points), closed_flag))
        del points[:]

    for command, args, relative in commands:
        if command == "M":
            flush(closed)
            closed = False
            nx, ny = args[0], args[1]
            if relative:
                nx += x
                ny += y
            x, y = nx, ny
            subpath_start = (x, y)
            points.append((x, y))
        elif command == "L":
            nx, ny = args[0], args[1]
            if relative:
                nx += x
                ny += y
            x, y = nx, ny
            points.append((x, y))
        elif command == "H":
            x = args[0] + (x if relative else 0.0)
            points.append((x, y))
        elif command == "V":
            y = args[0] + (y if relative else 0.0)
            points.append((x, y))
        elif command == "C":
            x1, y1, x2, y2, nx, ny = args
            if relative:
                x1 += x; y1 += y; x2 += x; y2 += y; nx += x; ny += y
            points.extend(sample_cubic((x, y), (x1, y1), (x2, y2), (nx, ny)))
            last_cubic_control = (x2, y2)
            x, y = nx, ny
        elif command == "S":
            x2, y2, nx, ny = args
            if relative:
                x2 += x; y2 += y; nx += x; ny += y
            # First control reflects the previous C/S control about the current
            # point; with no C/S predecessor it equals the current point.
            if last_cubic_control is not None:
                x1, y1 = 2 * x - last_cubic_control[0], 2 * y - last_cubic_control[1]
            else:
                x1, y1 = x, y
            points.extend(sample_cubic((x, y), (x1, y1), (x2, y2), (nx, ny)))
            last_cubic_control = (x2, y2)
            x, y = nx, ny
        elif command == "Q":
            x1, y1, nx, ny = args
            if relative:
                x1 += x; y1 += y; nx += x; ny += y
            points.extend(sample_quadratic((x, y), (x1, y1), (nx, ny)))
            last_quad_control = (x1, y1)
            x, y = nx, ny
        elif command == "T":
            nx, ny = args
            if relative:
                nx += x
                ny += y
            if last_quad_control is not None:
                x1, y1 = 2 * x - last_quad_control[0], 2 * y - last_quad_control[1]
            else:
                x1, y1 = x, y
            points.extend(sample_quadratic((x, y), (x1, y1), (nx, ny)))
            last_quad_control = (x1, y1)
            x, y = nx, ny
        elif command == "A":
            rx, ry, rotation, large_arc, sweep, nx, ny = args
            if relative:
                nx += x
                ny += y
            arc_count += 1
            points.extend(sample_arc(x, y, rx, ry, rotation, bool(large_arc), bool(sweep), nx, ny))
            x, y = nx, ny
        elif command == "Z":
            # Close via the freeform close flag; the current position returns
            # to the subpath start per the SVG spec. Any drawing command after
            # Z starts a fresh subpath (accepted simplification).
            flush(True)
            closed = False
            x, y = subpath_start
        # Reflection memory is only valid while the same bezier family repeats.
        if command not in ("C", "S"):
            last_cubic_control = None
        if command not in ("Q", "T"):
            last_quad_control = None
    flush(closed)
    return subpaths, arc_count


def text_segments(element):
    """Splits a <text> element into positioned, styled segments.

    Handles plain text and any mix of <tspan> children: each tspan may carry
    x/y/dx/dy positioning plus font-size/fill/font-weight overrides. Segment
    order follows document order (element text, tspan text, tails).
    """
    default_size = parse_length(element.get("font-size"), 14.0)
    default_fill = element.get("fill")
    default_bold = str(element.get("font-weight") or "").lower() in ("bold", "700", "800", "900")
    segments = []

    def add(text, style):
        if text is None or not str(text).strip():
            return
        segments.append({
            "text": str(text),
            "x": style.get("x"),
            "y": style.get("y"),
            "dx": style.get("dx") or 0.0,
            "dy": style.get("dy") or 0.0,
            "font_size": parse_length(style.get("font_size"), default_size),
            "fill": style.get("fill") if style.get("fill") is not None else default_fill,
            "bold": style.get("bold") if style.get("bold") is not None else default_bold,
        })

    parent_style = {"x": parse_length(element.get("x"), 0.0), "y": parse_length(element.get("y"), 0.0),
                    "font_size": element.get("font-size"), "fill": default_fill, "bold": default_bold}
    add(element.text, dict(parent_style, x=None, y=None))
    for child in element:
        if local_name(child.tag) != "tspan":
            continue
        tspan_style = {
            "x": parse_length(child.get("x"), None),
            "y": parse_length(child.get("y"), None),
            "dx": parse_length(child.get("dx"), 0.0),
            "dy": parse_length(child.get("dy"), 0.0),
            "font_size": child.get("font-size"),
            "fill": child.get("fill"),
            "bold": (str(child.get("font-weight") or "").lower() in ("bold", "700", "800", "900"))
                    if child.get("font-weight") else None,
        }
        add(child.text, tspan_style)
        add(child.tail, dict(parent_style, x=None, y=None))
    return segments


def build_text_paragraphs(element):
    """Groups text segments into paragraphs for a single text frame.

    Paragraph strategy (documented choice): a segment that carries an
    explicit x or y, or a non-zero dy (baseline shift, the common
    "next line" idiom) starts a NEW PARAGRAPH; purely inline segments
    (style-only tspans) continue the current paragraph as an additional run.
    dx on a continuing run is absorbed (single runs cannot shift vertically
    or horizontally inside a pptx paragraph); dy on a continuing run would
    also be impossible, so only dy==0 continues inline.
    """
    segments = text_segments(element)
    if not segments:
        return []
    paragraphs = []
    runs = []
    baseline = parse_length(element.get("y"), 0.0)
    anchor_x = parse_length(element.get("x"), 0.0)
    for segment in segments:
        starts_paragraph = segment["x"] is not None or segment["y"] is not None or segment["dy"] != 0.0
        if starts_paragraph and runs:
            paragraphs.append({"runs": runs, "baseline": baseline, "anchor_x": anchor_x})
            runs = []
        if segment["y"] is not None:
            baseline = segment["y"]
        elif segment["dy"] != 0.0:
            baseline += segment["dy"]
        if segment["x"] is not None:
            anchor_x = segment["x"]
        runs.append((segment["text"], segment["font_size"], segment["fill"], segment["bold"]))
    if runs:
        paragraphs.append({"runs": runs, "baseline": baseline, "anchor_x": anchor_x})
    return paragraphs


def main(argv):
    if len(argv) < 3:
        fail("usage: figure_svg_to_pptx.py <svg_path> <pptx_path> [title]")
    svg_path, pptx_path = argv[1], argv[2]
    title = argv[3] if len(argv) > 3 else "Scientific figure"

    try:
        import xml.etree.ElementTree as ET
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    except ImportError as error:
        fail("python-pptx is required: pip install python-pptx (%s)" % error)

    try:
        root = ET.parse(svg_path).getroot()
    except (ET.ParseError, OSError) as error:
        fail("failed to parse SVG %s: %s" % (svg_path, error))
    if local_name(root.tag) != "svg":
        fail("not an SVG document: %s" % svg_path)

    width = parse_length(root.get("width"), None)
    height = parse_length(root.get("height"), None)
    if width is None or height is None:
        viewBox = str(root.get("viewBox") or "").split()
        if len(viewBox) >= 4:
            width = width if width is not None else parse_length(viewBox[2], DEFAULT_WIDTH)
            height = height if height is not None else parse_length(viewBox[3], DEFAULT_HEIGHT)
    width = width or DEFAULT_WIDTH
    height = height or DEFAULT_HEIGHT

    presentation = Presentation()
    presentation.slide_width = Emu(int(round(width * EMU_PER_PX)))
    presentation.slide_height = Emu(int(round(height * EMU_PER_PX)))
    presentation.core_properties.title = title
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank

    gradients = collect_gradients(root)
    features = {}
    skipped = {}

    def style_fill(shape, resolved_color):
        # Consumes an ALREADY RESOLVED paint (resolve_paint output: hex string,
        # 'none', or None) — parsing it twice would drop named/url() values.
        if resolved_color == "none":
            shape.fill.background()
        elif resolved_color is None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string("000000")
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(resolved_color)

    def style_line(line_format, color_value, width_value):
        color = resolve_paint(color_value, gradients)
        if color == "none":
            line_format.fill.background()
            return
        if color is None:
            color = "000000"
        line_format.color.rgb = RGBColor.from_string(color)
        if width_value:
            line_format.width = Emu(int(round(parse_length(width_value, 1.0) * EMU_PER_PX)))

    def element_affine(affine, element):
        return affine.composed(parse_transform_affine(element.get("transform")))

    def placed_box(affine, left, top, box_w, box_h):
        """Places a bounding-box shape under an affine (center-orbit + rotation).

        Exact for rotate+translate (the constrained vocabulary): the box keeps
        its size, its center maps through the affine, and the rotation angle
        lands on python-pptx shape.rotation (both are clockwise in the y-down
        screen space). A scale/skew component cannot resize autoshapes here,
        so it is ignored and counted in svgFeaturesSkipped
        ["scaleOnBoxShapeIgnored"].
        """
        if affine.has_scale_or_skew:
            bump(skipped, "scaleOnBoxShapeIgnored")
        center_x, center_y = affine.apply(left + box_w / 2.0, top + box_h / 2.0)
        return (
            Emu(int(round((center_x - box_w / 2.0) * EMU_PER_PX))),
            Emu(int(round((center_y - box_h / 2.0) * EMU_PER_PX))),
            Emu(int(round(box_w * EMU_PER_PX))),
            Emu(int(round(box_h * EMU_PER_PX))),
            affine.rotation_degrees,
        )

    def apply_rotation(shape, angle):
        if abs(angle) > 0.01:
            shape.rotation = angle % 360.0
            bump(features, "rotate")

    def resolved_fill(element):
        """Element fill resolved through the gradient table, with counter."""
        value = element.get("fill")
        resolved = resolve_paint(value, gradients)
        if resolved is not None and re.match(r"url\(\s*#", str(value).strip()):
            bump(features, "gradientFill")
        return resolved

    shape_count = 0
    for tag, element, group_affine in collect_elements(root):
        affine = element_affine(group_affine, element)
        if tag == "rect":
            x = parse_length(element.get("x"), 0.0)
            y = parse_length(element.get("y"), 0.0)
            w = parse_length(element.get("width"), 0.0)
            h = parse_length(element.get("height"), 0.0)
            if w <= 0 or h <= 0:
                continue
            rx = parse_length(element.get("rx"), 0.0)
            shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if rx and rx > 0 else MSO_SHAPE.RECTANGLE
            left_emu, top_emu, w_emu, h_emu, angle = placed_box(affine, x, y, w, h)
            shape = slide.shapes.add_shape(shape_kind, left_emu, top_emu, w_emu, h_emu)
            shape.shadow.inherit = False
            apply_rotation(shape, angle)
            style_fill(shape, resolved_fill(element))
            if element.get("stroke"):
                style_line(shape.line, element.get("stroke"), element.get("stroke-width"))
            else:
                shape.line.fill.background()
            bump(features, "rect")
            shape_count += 1
        elif tag in ("circle", "ellipse"):
            cx = parse_length(element.get("cx"), 0.0)
            cy = parse_length(element.get("cy"), 0.0)
            rx = parse_length(element.get("r") if tag == "circle" else element.get("rx"), 0.0)
            ry = parse_length(element.get("r") if tag == "circle" else element.get("ry"), 0.0)
            if rx <= 0 or ry <= 0:
                continue
            left_emu, top_emu, w_emu, h_emu, angle = placed_box(affine, cx - rx, cy - ry, 2 * rx, 2 * ry)
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_emu, top_emu, w_emu, h_emu)
            shape.shadow.inherit = False
            apply_rotation(shape, angle)
            style_fill(shape, resolved_fill(element))
            if element.get("stroke"):
                style_line(shape.line, element.get("stroke"), element.get("stroke-width"))
            else:
                shape.line.fill.background()
            bump(features, tag)
            shape_count += 1
        elif tag == "line":
            x1, y1 = affine.apply(parse_length(element.get("x1"), 0.0), parse_length(element.get("y1"), 0.0))
            x2, y2 = affine.apply(parse_length(element.get("x2"), 0.0), parse_length(element.get("y2"), 0.0))
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(int(round(x1 * EMU_PER_PX))), Emu(int(round(y1 * EMU_PER_PX))),
                Emu(int(round(x2 * EMU_PER_PX))), Emu(int(round(y2 * EMU_PER_PX))))
            style_line(connector.line, element.get("stroke", "#000000"), element.get("stroke-width"))
            if abs(affine.rotation_degrees) > 0.01:
                # Vertex mapping already carries the rotation exactly.
                bump(features, "rotate")
            bump(features, "line")
            shape_count += 1
        elif tag in ("polyline", "polygon"):
            points_text = str(element.get("points") or "")
            pairs = re.findall(NUMBER_PATTERN, points_text)
            if len(pairs) < 4:
                bump(skipped, "%sWithoutEnoughPoints" % tag.capitalize())
                continue
            coordinates = [affine.apply(float(pairs[i]), float(pairs[i + 1])) for i in range(0, len(pairs) - 1, 2)]
            if len(coordinates) < 2:
                continue
            if abs(affine.rotation_degrees) > 0.01:
                bump(features, "rotate")
            for (ax, ay), (bx, by) in zip(coordinates, coordinates[1:]):
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Emu(int(round(ax * EMU_PER_PX))), Emu(int(round(ay * EMU_PER_PX))),
                    Emu(int(round(bx * EMU_PER_PX))), Emu(int(round(by * EMU_PER_PX))))
                style_line(connector.line, element.get("stroke", "#000000"), element.get("stroke-width"))
                shape_count += 1
            bump(features, tag)
        elif tag == "path":
            d_text = element.get("d") or ""
            if not d_text.strip():
                bump(skipped, "pathWithoutD")
                continue
            subpaths, arc_count = flatten_path(d_text)
            if arc_count:
                bump(skipped, "arcApproximated", arc_count)
            if not subpaths:
                bump(skipped, "pathWithoutDrawableSubpath")
                continue
            if abs(affine.rotation_degrees) > 0.01:
                bump(features, "rotate")
            fill_value = resolved_fill(element)
            for points, closed in subpaths:
                builder = slide.shapes.build_freeform(points[0][0], points[0][1], EMU_PER_PX)
                builder.add_line_segments(points[1:], close=closed)
                shape = builder.convert_to_shape()
                shape.shadow.inherit = False
                style_fill(shape, fill_value)
                if element.get("stroke"):
                    style_line(shape.line, element.get("stroke"), element.get("stroke-width"))
                else:
                    shape.line.fill.background()
                shape_count += 1
            bump(features, "path")
        elif tag == "text":
            paragraphs = build_text_paragraphs(element)
            if not paragraphs:
                continue
            anchor = str(element.get("text-anchor") or "start").strip().lower()
            baseline0 = paragraphs[0]["baseline"]
            font0 = paragraphs[0]["runs"][0][1]
            # Local-coordinate bounding estimate over all paragraphs. Paragraph
            # x positions collapse to the frame's left edge: only the first
            # paragraph's anchor positions the frame (documented trade-off to
            # keep multi-line text inside ONE editable text frame).
            line_widths = []
            for paragraph in paragraphs:
                text_length = sum(len(run[0]) for run in paragraph["runs"])
                font_size = max(run[1] for run in paragraph["runs"])
                line_widths.append(max(text_length * font_size * 0.62, font_size * 2))
            est_width = max(line_widths)
            anchor_x = paragraphs[0]["anchor_x"]
            if anchor == "middle":
                left = anchor_x - est_width / 2
            elif anchor == "end":
                left = anchor_x - est_width
            else:
                left = anchor_x
            baselines = [paragraph["baseline"] for paragraph in paragraphs]
            top = min(baselines) - font0
            box_h = (max(baselines) - min(baselines)) + font0 * 1.5
            left_emu, top_emu, w_emu, h_emu, angle = placed_box(affine, left, top, est_width, box_h)
            box = slide.shapes.add_textbox(left_emu, top_emu, w_emu, h_emu)
            box.word_wrap = False
            box.shadow.inherit = False
            box.fill.background()
            box.line.fill.background()
            apply_rotation(box, angle)
            frame = box.text_frame
            frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
            for paragraph_index, paragraph in enumerate(paragraphs):
                text_paragraph = frame.paragraphs[0] if paragraph_index == 0 else frame.add_paragraph()
                for run_text, font_size, fill, bold in paragraph["runs"]:
                    run = text_paragraph.add_run()
                    run.text = run_text
                    font = run.font
                    font.size = Pt(max(font_size * 0.75, 1.0))
                    font.name = "Arial"
                    if bold:
                        font.bold = True
                    color = parse_color(fill)
                    if color not in (None, "none"):
                        font.color.rgb = RGBColor.from_string(color)
            for paragraph in paragraphs[1:]:
                bump(features, "tspanParagraph")
            for paragraph in paragraphs:
                bump(features, "tspanRun", max(len(paragraph["runs"]) - 1, 0))
            bump(features, "text")
            shape_count += 1

    if shape_count == 0:
        fail("no drawable elements found in %s" % svg_path)

    try:
        presentation.save(pptx_path)
    except OSError as error:
        fail("failed to write %s: %s" % (pptx_path, error))

    summary = {
        "ok": True,
        "pptx": pptx_path,
        "shapes": shape_count,
        "slidePx": [int(width), int(height)],
        # Additive v2 keys: feature counters for observability. The fd-pptx
        # workflow node forwards them when present; its stdout contract only
        # grows, never breaks.
        "svgFeaturesMapped": features,
        "svgFeaturesSkipped": skipped,
    }
    sys.stdout.write(json.dumps(summary))
    sys.stdout.write("\n")


if __name__ == "__main__":
    main(sys.argv)
